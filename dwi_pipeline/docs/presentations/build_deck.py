#!/usr/bin/env python3
"""Build the DKT Connectome slide deck (PPTX, plus optional PNG previews).

    python3 build_deck.py            # writes preview/dkt_connectome/dkt_connectome.pptx
    python3 build_deck.py --png      # also writes preview/dkt_connectome/slide-NN.png
    python3 build_deck.py --figures  # regenerates the sample-subject figures

Slide content lives in SLIDES so both renderers stay in sync. One deck per
pipeline: give the next one its own DECK name and it lands in its own folder.

--figures reads the cohort results tree, which is not in version control; the
generated PNGs are committed so the deck builds without it.
"""

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
DECK = "dkt_connectome"
OUT_DIR = HERE / "preview" / DECK

# Sample subject shown on the results slide: a TBI participant with a manual
# lesion mask, so the run exercises inpainting and the disconnectome.
SAMPLE = "TBI011011"
SESSION = "ses-2WK"
TBI_ROOT = HERE.parents[1] / "dwi_test_TBI"
RESULTS = TBI_ROOT / f"sub-{SAMPLE}_fastsurfer_inpaint"
BIDS = TBI_ROOT / "bids"

# Annotation colour for the lesion outline — chosen to stay legible on greyscale
# MRI and to differ from every colour used elsewhere in the deck.
MARK = "#00C2FF"

W, H = 13.333, 7.5
LEFT = 0.95
BODY_W = W - 2 * LEFT

INK = "1A2332"
MUTED = "6B7280"
ACCENT = "1A5FB4"
RULE = "D0D5DD"

SANS = "Arial"
MONO = "Consolas"


SLIDES = [
    {
        "kind": "flow",
        "title": "DKT Connectome",
        "kicker": "BIDS App for lesion-aware structural connectomics",
        "lead": "BIDS  dwi/ + anat/T1w  (optional fmap/, lesion mask)   →   dkt_connectome.csv, 78 x 78",
        "steps": [
            ("Step 1", "QSIPrep", "Motion, eddy, distortion correction, T1w–DWI alignment", False),
            ("Step 1.5", "Inpaint", "Fill lesion on T1w before reconstruction", True),
            ("Step 2", "Recon", "FreeSurfer or FastSurfer surfaces + DKT labels", False),
            ("Step 3", "QSIRecon", "SS3T-CSD fibre orientations, ACT-HSVS tractography", False),
            ("Step 4", "Connectome", "Warp labels to tractography grid, count streamlines", False),
            ("Step 4.5", "Disconnectome", "Quantify connectivity lost to the lesion", True),
            ("Step 5", "Node strength", "Graph metrics and ENIGMA-style clinical report", False),
        ],
        "link": "Greyed steps are optional — they run only when their inputs exist",
    },
    {
        "kind": "table",
        "title": "Defaults",
        "kicker": "Snakemake engine · ./run locally · submit.sh on HPC · overridable per cohort",
        "cols": ["Topic", "Default", "Why"],
        "widths": [2.6, 4.2, 4.6],
        "rows": [
            ["Parcellation", "DKT, 78 nodes", "Available from both recon tools"],
            ["Tractography", "ACT-HSVS", "Anatomical priors from surfaces"],
            ["Edge weight", "Streamline counts", "Simple, widely reported"],
            ["Distortion", "Measured fieldmap", "Explicit --syn or --no-sdc otherwise"],
            ["Lesion", "Auto-detected", "From BIDS lesion mask sidecar"],
        ],
        "link": "Study-agnostic — any BIDS DWI + T1w cohort runs",
    },
    {
        "kind": "columns",
        "title": "Optional steps",
        "kicker": "Each is a separate concern, skipped by default unless its input is present",
        "columns": [
            {
                "tag": "Step 1.5",
                "name": "Inpainting",
                "lines": [
                    "Recon tools assume healthy anatomy",
                    "Diffusion model fills the cavity on T1w",
                    "Auto when a lesion mask exists",
                    "Diffusion data left untouched",
                ],
            },
            {
                "tag": "Step 4.5",
                "name": "Disconnectome",
                "lines": [
                    "How much connectivity the lesion cut",
                    "Excise voxels or drop streamlines",
                    "Opt-in with --disconnection",
                    "Primary connectome unchanged",
                ],
            },
            {
                "tag": "Step 5",
                "name": "Node strength",
                "lines": [
                    "Per-region strength and asymmetry",
                    "ENIGMA figures and report.pdf",
                    "Runs once a connectome exists",
                    "Separate container and repo",
                ],
            },
        ],
        "link": "Optional does not mean secondary — these carry the lesion-aware claims",
    },
    {
        "kind": "results",
        "layout": "side",
        "title": "TBI participant — lesion handling",
        "kicker": "Manual mask, right frontal · neuroLIT inpainting before reconstruction",
        "images": [
            ("fig-inpainting.png", "Step 1.5 — axial T1w, lesion outlined and ringed in blue"),
        ],
        "stats": [
            ("18.2 mL", "lesion (core + oedema)"),
            ("71%", "of lesion in right hemisphere"),
            ("r = 0.996", "T1w intact outside lesion"),
            ("78 x 78", "DKT nodes, none empty"),
        ],
        "link": "Inpainting rewrites only the masked voxels — recon then runs on near-normal anatomy",
    },
    {
        "kind": "table",
        "row_h": 0.47,
        "title": "TBI participant — node strength where the lesion lands",
        "kicker": "DKT regions ranked by lesion load, with the node strength and asymmetry of each",
        "cols": ["Region", "Lesion L", "Lesion R", "L strength", "R strength", "Str AI", "Vol AI"],
        "widths": [2.7, 1.4, 1.4, 1.8, 1.8, 1.55, 1.55],
        "rows": [
            ["Lateral orbitofrontal", "9.8%", "21.8%", "116,411", "113,107", "+0.014", "−0.014"],
            ["Middle temporal", "—", "8.0%", "250,161", "275,536", "−0.048", "+0.043"],
            ["Pars orbitalis", "—", "5.8%", "61,889", "62,103", "−0.002", "−0.025"],
            ["Insula", "—", "5.1%", "109,650", "109,316", "+0.002", "+0.008"],
            ["Pars triangularis", "—", "4.9%", "91,530", "135,264", "−0.193", "−0.111"],
            ["Superior temporal", "—", "4.2%", "232,350", "262,319", "−0.061", "+0.014"],
        ],
        "notes": [
            "Lesion load = share of the region's voxels inside the mask. The lesion touches 16 of the 78 DKT regions; these are the six largest.",
            "Strength = streamlines on all edges touching the node. AI = (L − R)/(L + R), so a negative value means the right side is stronger.",
            "The most damaged region is hit on both sides, so its AI stays near zero — asymmetry on its own is not a lesion detector.",
        ],
        "link": "Edge level is where the lesion shows: 1,062 of 2,897 edges weakened, 2 lost entirely, 2.9% of streamline weight",
    },
    {
        "kind": "table",
        "title": "Links",
        "kicker": "Documentation, code, and upstream engines",
        "cols": ["Resource", "Where"],
        "widths": [3.6, 7.8],
        "rows": [
            ["Documentation", "dkt-connectome.readthedocs.io"],
            ["Pipeline code", "github.com/phindagijimana/dkt_connectome"],
            ["Node strength", "github.com/phindagijimana/dwi-AI"],
            ["Inpainting model", "github.com/Deep-MI/lit"],
            ["Upstream engines", "qsiprep.readthedocs.io  ·  qsirecon.readthedocs.io"],
            ["Contact", "phindagiji@gmail.com"],
        ],
        "link": "",
    },
]


def _rgb(hex_str):
    return RGBColor.from_string(hex_str)


def _text(slide, x, y, w, h, runs, size=18, color=INK, bold=False,
          font=SANS, align=PP_ALIGN.LEFT, line=1.25):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    lines = runs if isinstance(runs, list) else [runs]
    for i, content in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = line
        run = para.add_run()
        run.text = content
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = _rgb(color)
    return box


def _rule(slide, x, y, w, color=RULE, thickness=0.012):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(thickness))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _header(slide, spec):
    _text(slide, LEFT, 0.72, BODY_W, 0.7, spec["title"], size=30, bold=True)
    _rule(slide, LEFT, 1.46, 0.85, ACCENT, 0.035)
    if spec.get("kicker"):
        _text(slide, LEFT, 1.62, BODY_W, 0.35, spec["kicker"], size=13, color=MUTED)


def _footer(slide, spec, number):
    if spec.get("link"):
        _text(slide, LEFT, 6.72, BODY_W - 1.0, 0.35, spec["link"], size=12, color=MUTED)
    if number:
        _text(slide, W - LEFT - 0.6, 6.72, 0.6, 0.35, str(number),
              size=12, color=MUTED, align=PP_ALIGN.RIGHT)


def render_title(slide, spec):
    _text(slide, LEFT, 2.45, BODY_W, 1.1, spec["title"], size=46, bold=True)
    _rule(slide, LEFT, 3.72, 1.3, ACCENT, 0.04)
    _text(slide, LEFT, 3.98, BODY_W, 0.5, spec["subtitle"], size=22, color=ACCENT)
    _text(slide, LEFT, 4.62, BODY_W, 0.5, spec["note"], size=15, color=MUTED)


def render_bullets(slide, spec):
    y = 2.15
    for label, text in spec["bullets"]:
        _text(slide, LEFT, y + 0.04, 1.9, 0.4, label, size=14, bold=True, color=ACCENT)
        _text(slide, LEFT + 2.05, y, BODY_W - 2.05, 0.7, text, size=19)
        y += 0.92


def render_flow(slide, spec):
    y = 2.05
    if spec.get("lead"):
        _text(slide, LEFT, y, BODY_W, 0.4, spec["lead"], size=14, color=INK)
        _rule(slide, LEFT, y + 0.42, BODY_W)
        y += 0.68

    for num, name, desc, optional in spec["steps"]:
        color = MUTED if optional else ACCENT
        name_color = MUTED if optional else INK
        _text(slide, LEFT, y + 0.03, 1.0, 0.35, num, size=13, bold=True, color=color)
        _text(slide, LEFT + 1.15, y, 2.5, 0.4, name, size=17, bold=True, color=name_color)
        _text(slide, LEFT + 3.75, y + 0.02, BODY_W - 3.75, 0.4, desc, size=15, color=MUTED)
        y += 0.57


def render_columns(slide, spec):
    cols = spec["columns"]
    gutter = 0.5
    width = (BODY_W - gutter * (len(cols) - 1)) / len(cols)

    for i, col in enumerate(cols):
        x = LEFT + i * (width + gutter)
        _text(slide, x, 2.15, width, 0.3, col["tag"], size=12, bold=True, color=ACCENT)
        _text(slide, x, 2.52, width, 0.45, col["name"], size=21, bold=True)
        _rule(slide, x, 3.08, width, ACCENT, 0.02)

        y = 3.28
        for entry in col["lines"]:
            _text(slide, x, y, width, 0.8, entry, size=14, color=MUTED, line=1.2)
            y += 0.62


def _draw_table(slide, spec, total_w, size=16, row_h=0.52):
    """Draw a borderless table starting at the left margin; returns bottom y."""
    widths = spec["widths"]
    xs, x = [], LEFT
    for width in widths:
        xs.append(x)
        x += width

    y = 2.15
    for i, col in enumerate(spec["cols"]):
        _text(slide, xs[i], y, widths[i] - 0.2, 0.35, col, size=12, bold=True, color=MUTED)
    y += 0.42
    _rule(slide, LEFT, y, total_w)
    y += 0.22

    for row in spec["rows"]:
        for i, cell in enumerate(row):
            _text(slide, xs[i], y, widths[i] - 0.2, 0.5, cell,
                  size=size, bold=(i == 0), color=INK if i == 0 else MUTED)
        y += row_h
        _rule(slide, LEFT, y - 0.12, total_w, "EDF0F3")
    return y


def render_table(slide, spec):
    y = _draw_table(slide, spec, BODY_W, row_h=spec.get("row_h", 0.52))
    for note in spec.get("notes", []):
        _text(slide, LEFT, y + 0.16, BODY_W, 0.4, note, size=11, color=MUTED)
        y += 0.30


IMG_TOP = 2.15
IMG_MAX_H = 2.8
IMG_GUTTER = 0.4
STATS_Y = 5.55


def _aspect(path):
    import matplotlib.image as mpimg

    h_px, w_px = mpimg.imread(path).shape[:2]
    return w_px / h_px


def _image_row(specs):
    """Lay images out left to right at a common height, filling the body width."""
    images = [(OUT_DIR / name, cap) for name, cap in specs]
    images = [(p, c) for p, c in images if p.exists()]
    if not images:
        return [], 0.0

    aspects = [_aspect(p) for p, _ in images]
    usable = BODY_W - IMG_GUTTER * (len(images) - 1)
    height = min(IMG_MAX_H, usable / sum(aspects))

    placed, x = [], LEFT
    for (path, caption), aspect in zip(images, aspects):
        w = height * aspect
        placed.append((path, caption, x, w))
        x += w + IMG_GUTTER
    return placed, height


def render_results(slide, spec):
    if spec.get("layout") == "side":
        y = 2.25
        for value, label in spec["stats"]:
            _text(slide, LEFT, y, 3.2, 0.4, value, size=20, bold=True)
            _text(slide, LEFT, y + 0.36, 3.2, 0.5, label, size=11.5, color=MUTED)
            y += 0.95

        name, caption = spec["images"][0]
        path = OUT_DIR / name
        if not path.exists():
            return
        x = LEFT + 3.5
        max_w = W - LEFT - x
        height = min(4.1, max_w / _aspect(path))
        width = height * _aspect(path)
        slide.shapes.add_picture(str(path), Inches(x), Inches(IMG_TOP),
                                 width=Inches(width), height=Inches(height))
        _text(slide, x, IMG_TOP + height + 0.14, max_w, 0.3, caption, size=11, color=MUTED)
        return

    placed, height = _image_row(spec["images"])
    for path, caption, x, w in placed:
        slide.shapes.add_picture(str(path), Inches(x), Inches(IMG_TOP),
                                 width=Inches(w), height=Inches(height))
        _text(slide, x, IMG_TOP + height + 0.14, w, 0.3, caption, size=11, color=MUTED)

    stats = spec["stats"]
    col = BODY_W / len(stats)
    for i, (value, label) in enumerate(stats):
        x = LEFT + i * col
        _text(slide, x, STATS_Y, col - 0.3, 0.4, value, size=20, bold=True)
        _text(slide, x, STATS_Y + 0.36, col - 0.3, 0.4, label, size=11.5, color=MUTED)


RENDERERS = {
    "title": render_title,
    "bullets": render_bullets,
    "flow": render_flow,
    "columns": render_columns,
    "results": render_results,
    "table": render_table,
}


def _lesion_panels():
    """Axial T1w before and after inpainting, with the lesion ringed on both."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import nibabel as nib
    import numpy as np
    from matplotlib.patches import Ellipse

    inp = RESULTS / "inpainted" / f"sub-{SAMPLE}" / SESSION
    before_f = BIDS / f"sub-{SAMPLE}" / SESSION / "anat" / f"sub-{SAMPLE}_{SESSION}_acq-3D_T1w.nii.gz"
    after_f = inp / "inpainting_volumes" / "inpainting_result.nii.gz"
    mask_f = inp / "lesion_mask_prepared.nii.gz"
    if not (before_f.exists() and after_f.exists() and mask_f.exists()):
        return None

    before = nib.load(before_f).get_fdata()
    after = nib.load(after_f).get_fdata()
    mask = nib.load(mask_f).get_fdata() > 0

    # Axial slice through the largest cross-section of the lesion.
    z = int(mask.sum(axis=(0, 1)).argmax())

    def plane(vol):
        return vol[:, :, z].T

    m = plane(mask)
    # Crop to the head so the panels are not mostly background.
    head = plane(before) > plane(before).max() * 0.06
    rows, cols = np.where(head)
    r0, r1 = max(rows.min() - 4, 0), min(rows.max() + 5, head.shape[0])
    c0, c1 = max(cols.min() - 4, 0), min(cols.max() + 5, head.shape[1])

    ys, xs = np.where(m)
    ring_xy = ((xs.min() + xs.max()) / 2 - c0, (ys.min() + ys.max()) / 2 - r0)
    ring_w = (xs.max() - xs.min()) * 1.3 + 10
    ring_h = (ys.max() - ys.min()) * 1.3 + 10

    fig, axes = plt.subplots(1, 2, figsize=(6.6, 3.5), dpi=200)
    for ax, vol, title in ((axes[0], before, "Before — lesion"),
                           (axes[1], after, "After — inpainted")):
        img = plane(vol)
        ax.imshow(img[r0:r1, c0:c1], cmap="gray", origin="lower",
                  vmin=0, vmax=np.percentile(img[img > 0], 99.5))
        ax.contour(m[r0:r1, c0:c1], levels=[0.5], colors=MARK, linewidths=1.1)
        ax.add_patch(Ellipse(ring_xy, ring_w, ring_h, fill=False,
                             edgecolor=MARK, linewidth=1.8, linestyle="--"))
        ax.set_title(title, fontsize=11, color="#1A2332")
        ax.axis("off")

    fig.tight_layout(pad=0.3)
    path = OUT_DIR / "fig-inpainting.png"
    fig.savefig(path, facecolor="white")
    plt.close(fig)
    return path


def build_figures():
    """Render sample-subject figures from the cohort results tree."""
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    written = []

    path = _lesion_panels()
    if path:
        written.append(path)

    return written


def build_pptx(out_path):
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W), Inches(H)
    blank = prs.slide_layouts[6]

    for i, spec in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank)
        if spec["kind"] != "title":
            _header(slide, spec)
            _footer(slide, spec, i)
        else:
            _footer(slide, spec, None)
        RENDERERS[spec["kind"]](slide, spec)

    prs.save(out_path)
    return out_path


def build_png(out_dir):
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    out_dir.mkdir(parents=True, exist_ok=True)
    hexc = lambda c: "#" + c

    def put(fig, x, y, text, size, color, bold=False, mono=False):
        fig.text(
            x / W, 1 - y / H, text,
            fontsize=size, color=hexc(color), va="top", ha="left",
            fontweight="bold" if bold else "normal",
            family="monospace" if mono else "sans-serif",
        )

    def line(fig, x, y, w, color, thick=1.2):
        fig.add_artist(plt.Line2D(
            [x / W, (x + w) / W], [1 - y / H, 1 - y / H],
            color=hexc(color), linewidth=thick,
        ))

    paths = []
    for i, spec in enumerate(SLIDES, start=1):
        fig = plt.figure(figsize=(W, H), dpi=110)
        fig.patch.set_facecolor("white")

        if spec["kind"] == "title":
            put(fig, LEFT, 2.45, spec["title"], 40, INK, bold=True)
            line(fig, LEFT, 3.80, 1.3, ACCENT, 3)
            put(fig, LEFT, 4.02, spec["subtitle"], 20, ACCENT)
            put(fig, LEFT, 4.66, spec["note"], 13, MUTED)
        else:
            put(fig, LEFT, 0.78, spec["title"], 27, INK, bold=True)
            line(fig, LEFT, 1.50, 0.85, ACCENT, 3)
            if spec.get("kicker"):
                put(fig, LEFT, 1.66, spec["kicker"], 11.5, MUTED)

            if spec["kind"] == "bullets":
                y = 2.20
                for label, text in spec["bullets"]:
                    put(fig, LEFT, y + 0.03, label, 12.5, ACCENT, bold=True)
                    put(fig, LEFT + 2.05, y, text, 17, INK)
                    y += 0.92

            elif spec["kind"] == "flow":
                y = 2.10
                if spec.get("lead"):
                    put(fig, LEFT, y, spec["lead"], 12.5, INK)
                    line(fig, LEFT, y + 0.42, BODY_W, RULE)
                    y += 0.68
                for num, name, desc, optional in spec["steps"]:
                    col = MUTED if optional else ACCENT
                    ncol = MUTED if optional else INK
                    put(fig, LEFT, y + 0.02, num, 11.5, col, bold=True)
                    put(fig, LEFT + 1.15, y, name, 15, ncol, bold=True)
                    put(fig, LEFT + 3.75, y + 0.02, desc, 13, MUTED)
                    y += 0.57

            elif spec["kind"] == "columns":
                cols = spec["columns"]
                gutter = 0.5
                cw = (BODY_W - gutter * (len(cols) - 1)) / len(cols)
                for j, col in enumerate(cols):
                    x = LEFT + j * (cw + gutter)
                    put(fig, x, 2.20, col["tag"], 10.5, ACCENT, bold=True)
                    put(fig, x, 2.57, col["name"], 19, INK, bold=True)
                    line(fig, x, 3.14, cw, ACCENT, 2)
                    y = 3.34
                    for entry in col["lines"]:
                        put(fig, x, y, entry, 12.5, MUTED)
                        y += 0.62

            elif spec["kind"] == "table":
                widths = spec["widths"]
                xs, x = [], LEFT
                for width in widths:
                    xs.append(x)
                    x += width
                y = 2.20
                for j, col in enumerate(spec["cols"]):
                    put(fig, xs[j], y, col, 10.5, MUTED, bold=True)
                y += 0.42
                line(fig, LEFT, y, BODY_W, RULE)
                y += 0.28
                row_h = spec.get("row_h", 0.52)
                for row in spec["rows"]:
                    for j, cell in enumerate(row):
                        put(fig, xs[j], y, cell, 14, INK if j == 0 else MUTED, bold=(j == 0))
                    y += row_h
                    line(fig, LEFT, y - 0.12, BODY_W, "EDF0F3")
                for note in spec.get("notes", []):
                    put(fig, LEFT, y + 0.18, note, 10, MUTED)
                    y += 0.30

            elif spec["kind"] == "results" and spec.get("layout") == "side":
                y = 2.25
                for value, label in spec["stats"]:
                    put(fig, LEFT, y, value, 18, INK, bold=True)
                    put(fig, LEFT, y + 0.38, label, 10.5, MUTED)
                    y += 0.95

                name, caption = spec["images"][0]
                path = OUT_DIR / name
                if path.exists():
                    x = LEFT + 3.5
                    max_w = W - LEFT - x
                    ih = min(4.1, max_w / _aspect(path))
                    iw = ih * _aspect(path)
                    ax = fig.add_axes([x / W, 1 - (IMG_TOP + ih) / H, iw / W, ih / H])
                    ax.imshow(plt.imread(path))
                    ax.axis("off")
                    put(fig, x, IMG_TOP + ih + 0.16, caption, 9.5, MUTED)

            elif spec["kind"] == "results":
                placed, height = _image_row(spec["images"])
                for path, caption, x, w in placed:
                    ax = fig.add_axes([x / W, 1 - (IMG_TOP + height) / H, w / W, height / H])
                    ax.imshow(plt.imread(path))
                    ax.axis("off")
                    put(fig, x, IMG_TOP + height + 0.16, caption, 9.5, MUTED)

                stats = spec["stats"]
                cw = BODY_W / len(stats)
                for j, (value, label) in enumerate(stats):
                    x = LEFT + j * cw
                    put(fig, x, STATS_Y, value, 18, INK, bold=True)
                    put(fig, x, STATS_Y + 0.38, label, 10.5, MUTED)

            if spec.get("link"):
                put(fig, LEFT, 6.78, spec["link"], 10.5, MUTED)
            fig.text((W - LEFT) / W, 1 - 6.78 / H, str(i),
                     fontsize=10.5, color=hexc(MUTED), va="top", ha="right")

        path = out_dir / f"slide-{i:02d}.png"
        fig.savefig(path, facecolor="white")
        plt.close(fig)
        paths.append(path)
    return paths


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--png", action="store_true", help="also render preview images")
    ap.add_argument("--figures", action="store_true",
                    help="regenerate sample-subject figures from the results tree")
    args = ap.parse_args()

    if args.figures:
        written = build_figures()
        if not written:
            print(f"no results found under {RESULTS} — figures unchanged")
        for path in written:
            print(f"wrote {path}")

    pptx_path = build_pptx(OUT_DIR / f"{DECK}.pptx")
    print(f"wrote {pptx_path}  ({len(SLIDES)} slides)")

    if args.png:
        for path in build_png(OUT_DIR):
            print(f"wrote {path}")


if __name__ == "__main__":
    main()
