#!/usr/bin/env python3
"""Build the AutoHS slide deck (PPTX, plus optional PNG previews).

    python3 build_deck_autohs.py            # writes preview/AutoHS/AutoHS.pptx
    python3 build_deck_autohs.py --png      # also writes preview/AutoHS/slide-NN.png
    python3 build_deck_autohs.py --figures  # regenerates the sample-subject figure

Sibling of build_deck.py — same layout helpers, different SLIDES + sample. The
sample subject is job 29b52ad1 (sub-001_ses-1, IDEAS BIDS sample), for which
AutoHS has already produced a full segmentation + hippocampal report.
"""

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
DECK = "AutoHS"
OUT_DIR = HERE / "preview" / DECK

# Sample AutoHS run: sub-001 from the IDEAS BIDS sample. Job 29b52ad1 has a
# full FreeSurfer subject_dir + report.json so we can draw the coronal T1w
# slice and pull real hippocampus volumes.
AUTOHS_ROOT = Path("/mnt/nfs/home/urmc-sh.rochester.edu/pndagiji/Documents/AutoHS")
SAMPLE_JOB = "29b52ad1"
SAMPLE_ID = "sub-001"
JOB_ROOT = AUTOHS_ROOT / "data" / "jobs" / SAMPLE_JOB
REPORT_JSON = JOB_ROOT / "output" / "report.json"
T1W = JOB_ROOT / "input" / "sub-001_ses-1_T1w.nii.gz"
FS_DIR = JOB_ROOT / "freesurfer" / f"job_{SAMPLE_JOB}"

# Colour to outline the hippocampi — same accent used elsewhere but shifted
# to keep left / right visually distinct.
HIPPO_L = "#00C2FF"
HIPPO_R = "#FF6B35"

# ---- Layout constants (mirrors build_deck.py) --------------------------------
W, H = 13.333, 7.5
LEFT = 0.95
BODY_W = W - 2 * LEFT

INK = "1A2332"
MUTED = "6B7280"
ACCENT = "1A5FB4"
RULE = "D0D5DD"

SANS = "Arial"
MONO = "Consolas"


# ---- Slide content -----------------------------------------------------------

def _load_metrics():
    """Pull real hippocampus stats for sub-001 from the AutoHS report.json."""
    if not REPORT_JSON.exists():
        return {}
    try:
        d = json.loads(REPORT_JSON.read_text())
        return d.get("metrics", {})
    except Exception:
        return {}


def _slides():
    m = _load_metrics()
    left = m.get("left_hippocampus_mm3", 5212.31)
    right = m.get("right_hippocampus_mm3", 5549.72)
    ai = m.get("asymmetry_index", -0.0314)
    laterality = m.get("laterality", "Symmetric")
    lat_thr = m.get("laterality_threshold", 0.05)
    left_thr = m.get("left_hs_threshold", -0.0708)
    right_thr = m.get("right_hs_threshold", 0.0469)
    hs_class = m.get("hs_classification", "Balanced (No HS)").split("\n")[0]

    return [
        {
            "kind": "flow",
            "title": "AutoHS",
            "kicker": "BIDS App for automated hippocampal sclerosis screening from T1w MRI",
            "lead": "BIDS  anat/T1w   →   per-subject HS report (PDF + JSON), BIDS-compliant derivatives",
            "steps": [
                ("Step 1", "Segment",
                 "FreeSurfer or FastSurfer — cortical parcellation + subcortical aseg", False),
                ("Step 2", "Extract volumes",
                 "Left- and right-hippocampus mm³ from FreeSurfer aseg", False),
                ("Step 3", "Asymmetry index",
                 "AI = (L − R) / (L + R)   —   direction and magnitude of hippocampal asymmetry", False),
                ("Step 4", "Classify",
                 "Compare AI against published HS thresholds (IDEAS-derived, per-side)", False),
                ("Step 5", "Report",
                 "PDF + JSON, BIDS derivative layout, ready for review or dashboard", False),
            ],
            "link": "MIT-licensed BIDS App · runs Docker or Apptainer · single T1w per subject",
        },
        {
            "kind": "table",
            "title": "Defaults",
            "kicker": "Threshold-based classifier, no model to train, no reference cohort required at inference",
            "cols": ["Topic", "Default", "Why"],
            "widths": [2.8, 4.6, 4.0],
            "rows": [
                ["Segmentation", "FastSurfer (default) or FreeSurfer 7.4.1",
                 "Deep-learning speed with recon-all option"],
                ["Volume source", "aseg — Left-/Right-Hippocampus (labels 17, 53)",
                 "Robust, well-validated in vivo hippocampus"],
                ["Asymmetry index", "AI = (L − R) / (L + R)",
                 "Standard laterality index, dimensionless"],
                ["Laterality threshold", f"±{lat_thr}",
                 "Reports Left/Right/Symmetric for the AI report"],
                ["HS thresholds", f"Left HS: AI < {left_thr:.4f}   ·   Right HS: AI > {right_thr:.4f}",
                 "Derived on IDEAS cohort — validated on epilepsy surgical patients"],
                ["Output", "report.json + report.pdf + summary.txt",
                 "Machine-readable metrics plus a clinical PDF"],
            ],
            "link": "One knob per user: --fastsurfer (default) or --freesurfer",
        },
        {
            "kind": "results",
            "layout": "side",
            "title": f"Sample subject — {SAMPLE_ID}, IDEAS BIDS release",
            "kicker": "T1w with left (blue) and right (orange) hippocampus overlays · AutoHS auto-classified as No HS",
            "images": [
                ("fig-hippocampus.png", "Coronal T1w slice through the largest hippocampal cross-section; overlay from aparc+aseg.mgz"),
            ],
            "stats": [
                (f"{int(round(left))} mm³", "left hippocampus"),
                (f"{int(round(right))} mm³", "right hippocampus"),
                (f"AI = {ai:+.3f}", f"laterality: {laterality}"),
                (hs_class, "HS classification"),
            ],
            "link": (
                f"|AI|={abs(ai):.3f} is inside both HS thresholds "
                f"(L={abs(left_thr):.3f}, R={abs(right_thr):.3f}) — subject flagged as balanced/no HS"
            ),
        },
    ]


# ---- Rendering primitives — identical to build_deck.py -----------------------

def _rgb(hex_str):
    return RGBColor.from_string(hex_str)


def _text(slide, x, y, w, h, runs, size=18, color=INK, bold=False,
          font=SANS, align=PP_ALIGN.LEFT, line=1.25):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    lines = runs if isinstance(runs, list) else [runs]
    for i, content in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = line
        run = para.add_run()
        run.text = content
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = _rgb(color)
    return box


def _rule(slide, x, y, w, color=RULE, thickness=0.012):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(thickness))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _header(slide, spec):
    _text(slide, LEFT, 0.72, BODY_W, 0.7, spec["title"], size=30, bold=True)
    _rule(slide, LEFT, 1.46, 0.85, ACCENT, 0.035)
    if spec.get("kicker"):
        _text(slide, LEFT, 1.62, BODY_W, 0.35, spec["kicker"], size=13, color=MUTED)


def _footer(slide, spec, number):
    if spec.get("link"):
        _text(slide, LEFT, 6.72, BODY_W - 1.0, 0.35, spec["link"], size=12, color=MUTED)
    if number:
        _text(slide, W - LEFT - 0.6, 6.72, 0.6, 0.35, str(number),
              size=12, color=MUTED, align=PP_ALIGN.RIGHT)


def render_flow(slide, spec):
    y = 2.05
    if spec.get("lead"):
        _text(slide, LEFT, y, BODY_W, 0.4, spec["lead"], size=14, color=INK)
        _rule(slide, LEFT, y + 0.42, BODY_W)
        y += 0.68

    for num, name, desc, optional in spec["steps"]:
        color = MUTED if optional else ACCENT
        name_color = MUTED if optional else INK
        _text(slide, LEFT, y + 0.03, 1.0, 0.35, num, size=13, bold=True, color=color)
        _text(slide, LEFT + 1.15, y, 2.5, 0.4, name, size=17, bold=True, color=name_color)
        _text(slide, LEFT + 3.75, y + 0.02, BODY_W - 3.75, 0.4, desc, size=15, color=MUTED)
        y += 0.57


def _draw_table(slide, spec, total_w, size=16, row_h=0.52):
    widths = spec["widths"]
    xs, x = [], LEFT
    for width in widths:
        xs.append(x)
        x += width

    y = 2.15
    for i, col in enumerate(spec["cols"]):
        _text(slide, xs[i], y, widths[i] - 0.2, 0.35, col, size=12, bold=True, color=MUTED)
    y += 0.42
    _rule(slide, LEFT, y, total_w)
    y += 0.22

    for row in spec["rows"]:
        for i, cell in enumerate(row):
            _text(slide, xs[i], y, widths[i] - 0.2, 0.5, cell,
                  size=size, bold=(i == 0), color=INK if i == 0 else MUTED)
        y += row_h
        _rule(slide, LEFT, y - 0.12, total_w, "EDF0F3")
    return y


def render_table(slide, spec):
    _draw_table(slide, spec, BODY_W, row_h=spec.get("row_h", 0.52))


IMG_TOP = 2.15


def _aspect(path):
    import matplotlib.image as mpimg
    h_px, w_px = mpimg.imread(path).shape[:2]
    return w_px / h_px


def render_results(slide, spec):
    y = 2.25
    for value, label in spec["stats"]:
        _text(slide, LEFT, y, 3.4, 0.4, value, size=18, bold=True)
        _text(slide, LEFT, y + 0.36, 3.4, 0.5, label, size=11.5, color=MUTED)
        y += 0.95

    name, caption = spec["images"][0]
    path = OUT_DIR / name
    if not path.exists():
        return
    x = LEFT + 3.7
    max_w = W - LEFT - x
    height = min(4.1, max_w / _aspect(path))
    width = height * _aspect(path)
    slide.shapes.add_picture(str(path), Inches(x), Inches(IMG_TOP),
                             width=Inches(width), height=Inches(height))
    _text(slide, x, IMG_TOP + height + 0.14, max_w, 0.3, caption, size=11, color=MUTED)


RENDERERS = {
    "flow": render_flow,
    "table": render_table,
    "results": render_results,
}


# ---- Sample-subject figure ---------------------------------------------------

def _hippocampus_panel():
    """Coronal T1w slice with L/R hippocampi outlined from aparc+aseg."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import nibabel as nib
    import numpy as np
    from scipy.ndimage import binary_dilation

    aseg_path = None
    for cand in [FS_DIR / "mri" / "aparc.DKTatlas+aseg.mgz",
                 FS_DIR / "mri" / "aparc+aseg.mgz",
                 FS_DIR / "mri" / "aparc.DKTatlas+aseg.deep.mgz",
                 FS_DIR / "mri" / "aseg.mgz"]:
        if cand.exists():
            aseg_path = cand
            break
    if aseg_path is None or not T1W.exists():
        return None

    t1 = nib.load(str(T1W))
    aseg = nib.load(str(aseg_path))

    # Resample aseg to the T1w grid if needed (FreeSurfer's conformed → orig).
    if aseg.shape != t1.shape:
        from nibabel.processing import resample_from_to
        aseg = resample_from_to(aseg, t1, order=0)

    t1_arr = np.asanyarray(t1.dataobj)
    seg_arr = np.rint(np.asanyarray(aseg.dataobj)).astype(int)

    left_mask = (seg_arr == 17)
    right_mask = (seg_arr == 53)
    both = left_mask | right_mask
    if not both.any():
        return None

    # Coronal slice through the biggest hippocampal cross-section.
    y = int(both.sum(axis=(0, 2)).argmax())

    t1_slice = np.rot90(t1_arr[:, y, :])
    left_slice = np.rot90(left_mask[:, y, :])
    right_slice = np.rot90(right_mask[:, y, :])

    # Contour = dilated mask minus mask; two-pixel outline for visibility.
    def _outline(m):
        d = binary_dilation(m, iterations=2)
        return d & ~m

    left_edge = _outline(left_slice)
    right_edge = _outline(right_slice)

    fig, ax = plt.subplots(figsize=(6.0, 5.6), dpi=200)
    lo, hi = np.percentile(t1_slice[t1_slice > 0], [1, 99])
    ax.imshow(t1_slice, cmap="gray", vmin=lo, vmax=hi, interpolation="none")

    overlay_l = np.zeros((*left_edge.shape, 4))
    overlay_l[left_edge] = [0.0, 0.76, 1.0, 1.0]   # HIPPO_L 00C2FF
    ax.imshow(overlay_l, interpolation="none")

    overlay_r = np.zeros((*right_edge.shape, 4))
    overlay_r[right_edge] = [1.0, 0.42, 0.21, 1.0]  # HIPPO_R FF6B35
    ax.imshow(overlay_r, interpolation="none")

    ax.set_xticks([]); ax.set_yticks([])
    for s in ax.spines.values():
        s.set_visible(False)
    fig.tight_layout(pad=0)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    path = OUT_DIR / "fig-hippocampus.png"
    fig.savefig(path, bbox_inches="tight", pad_inches=0.02, facecolor="white")
    plt.close(fig)
    return path


# ---- Build entry point -------------------------------------------------------

def build_pptx():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    slides = _slides()
    for idx, spec in enumerate(slides, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        _header(slide, spec)
        RENDERERS[spec["kind"]](slide, spec)
        _footer(slide, spec, idx)

    out = OUT_DIR / f"{DECK}.pptx"
    prs.save(str(out))
    return out, len(slides)


def build_pngs(n_slides):
    """Optional: render each slide to a PNG via LibreOffice."""
    import shutil
    import subprocess
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        print("(skipping PNGs — no LibreOffice)")
        return
    pptx = OUT_DIR / f"{DECK}.pptx"
    subprocess.run([soffice, "--headless", "--convert-to", "png",
                    "--outdir", str(OUT_DIR), str(pptx)], check=False)


def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--png", action="store_true", help="render slide-NN.png alongside the pptx")
    ap.add_argument("--figures", action="store_true",
                    help="regenerate the sample-subject figure from AutoHS outputs")
    args = ap.parse_args()

    if args.figures or not (OUT_DIR / "fig-hippocampus.png").exists():
        path = _hippocampus_panel()
        if path:
            print(f"figure -> {path}")
        else:
            print("WARNING: could not build hippocampus figure "
                  "(missing T1w, aseg, or the sample AutoHS job)")

    out, n = build_pptx()
    print(f"wrote {out} ({n} slides)")

    if args.png:
        build_pngs(n)


if __name__ == "__main__":
    main()
